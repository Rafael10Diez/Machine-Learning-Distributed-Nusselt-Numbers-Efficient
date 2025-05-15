# python C:\Users\rafae\Dropbox\workspace\code\202109_stag_interp_v5\classic_ML_use_Cf_Nu_search_v2\post\post_code\code_pptx\gen_pptx\mk_pptx.py
# ------------------------------------------------------------------------
#                  Generic Libraries
# ------------------------------------------------------------------------

from  datetime        import  datetime

from  os.path         import  dirname, basename, abspath
from  os.path         import  join  as  pjoin
from  os              import  walk  as  os_walk

from  PIL             import  Image

from  pptx            import  Presentation
from  pptx.util       import  Inches, Pt
from  pptx.enum.text  import  PP_ALIGN
from  copy            import  deepcopy

# ------------------------------------------------------------------------
#                  Basic Definitions
# ------------------------------------------------------------------------

file_dir       =  dirname(abspath(__file__))
tstamp         =  lambda   :  datetime.now().strftime("%Y%m%d_%H%M%S")
lfilter        =  lambda f,x: list(filter(f,x))
lmap           =  lambda f,x: list(map(f,x))
deepdirname    =  lambda x,n: x if n<1 else deepdirname(dirname(x),n-1)
post_output    =  pjoin(deepdirname(abspath(__file__),4),'post_output')

def reader(fname):
    with open(fname,'r') as f:
        return [x.rstrip('\n') for x in f]

def pop1(A):
    assert len(A) == 1
    return A[0]

def linux_path(x):
    cL  = r' / '.replace(' ','')
    cW  = r' \ '.replace(' ','')
    cL2 = cL*2
    for c in [cW,cL2]:
        while c in x:
            x = x.replace(c,cL)
    return x

# ------------------------------------------------------------------------
#                  Presentation Class
# ------------------------------------------------------------------------

class Mk_Presentation:
    def __init__(self, all_fname_pics, fname_pptx, D_qtypes):
        
        # register type of study
        self.all_fname_pics  =  all_fname_pics
        
        # create presentation object
        prs              =  Presentation(pjoin(file_dir,'template.pptx'))
        
        # define n_study
        # selfl.n_study     =  {'n11': 11, 'n6': 6, 'n23': 23}[type_study]
        
        for q in ['qmode']:
            
            q_type  =  D_qtypes[q]
            mk_tag  =  lambda b: ('('*b + {'Nu_classic': 'Nusselt Number', 
                                           'Qdot'      : 'Heat Transfer' ,
                                           'Cf_classic': 'Skin Friction' ,
                                           'Fx'        : 'Drag Forces'   }[q_type] + ')'*b)
            
            # make title slide (Fx/Qdot)
            tslide  =  prs.slides.add_slide(prs.slide_layouts[0])
            
            # initialize empty text fields
            for p in tslide.shapes:
                p.text  =  ' '
            
            # generate title text
            tslide.shapes.title.text  =  mk_tag(False)
            
            q_prefix = f'{q}_{q_type}'
            # add subsequent slides
            self.add_subset_slides(prs, mk_tag(True), q_prefix, all_fname_pics)
        
        self.fname_pptx  =  fname_pptx
        prs.save(self.fname_pptx)
    
    def add_subset_slides(self, prs, tag, q_prefix, all_fname_pics):
        for D in all_fname_pics:
            all_pics, all_text       =  self.mk_all_pics_text(D, q_prefix)
            slide                    =  prs.slides.add_slide(prs.slide_layouts[5])
            slide.shapes.title.text  =  f"Case {D['slide_name']}\n" + tag
            
            for pic,txt in zip(all_pics, all_text):
                
                img      =  slide.shapes.add_picture(           pic['fname' ]  ,
                                                                pic['left'  ]  ,
                                                                pic['top'   ]  ,
                                                     width   =  pic['width' ]  ,
                                                     height  =  pic['height']  )
                
                tf       =  slide.shapes.add_textbox( txt['left'  ]  ,
                                                      txt['top'   ]  ,
                                                      txt['width' ]  ,
                                                      txt['height']  ).text_frame
                tf.text  =   txt['label']
                for p in tf.paragraphs:
                    p.alignment  =  PP_ALIGN.CENTER
                    p.font.size  =  Pt(20)
    
    def mk_all_pics_text(self, D, q_prefix):
        
        all_pics   =  [ self.fix_info({'fname': D[f'{q_prefix}_dns']['fname_png']   , 'left': 1.11 , 'top': 1.89 , 'width': 6.15 , 'height': 3.07 }) ,
                        self.fix_info({'fname': D[ 'H']['fname_png']                , 'left': 8.74 , 'top': 1.89 , 'width': 6.15 , 'height': 3.07 }) ,
                        self.fix_info({'fname': D[f'{q_prefix}_pred']['fname_png']  , 'left': 1.11 , 'top': 5.64 , 'width': 6.15 , 'height': 3.07 }) ,
                        self.fix_info({'fname': D[f'{q_prefix}_error']['fname_png'] , 'left': 8.74 , 'top': 5.64 , 'width': 6.15 , 'height': 3.07 }) ]
        
        all_text   =  [ {'label': 'DNS Labels'        , 'left': Inches( 3.16) , 'top': Inches(1.43) , 'width': Inches(2.05) , 'height': Inches(0.4 ) } ,
                        {'label': 'Height Map'        , 'left': Inches(10.98) , 'top': Inches(1.43) , 'width': Inches(2.05) , 'height': Inches(0.4 ) } ,
                        {'label': 'Predicted (ML)'    , 'left': Inches( 3.16) , 'top': Inches(5.18) , 'width': Inches(2.05) , 'height': Inches(0.4 ) } ,
                        {'label': 'Errors Difference' , 'left': Inches(10.98) , 'top': Inches(5.18) , 'width': Inches(2.05) , 'height': Inches(0.4 ) } ]
        
        return all_pics, all_text
    
    def fix_info(self, D):
        D['fname'] = self.__fix_path(D['fname'])
        
        nw, nh       =  Image.open(D['fname']).size
        ideal_width  =  D['height']* nw / float(nh)
        
        w_center    =  D['left'] + 0.5 * D['width']
        
        # mk new_left, new_width
        D['left']   =  w_center - 0.5 * ideal_width
        D['width']  =  ideal_width
        
        for k in ('left','top','width','height'):
            D[k] = Inches(D[k])
        
        return D
    
    @staticmethod
    def __fix_path(x):
        p0    =  linux_path(post_output)
        _,p1  =  linux_path(x).split(basename(post_output))
        assert (p0[-1] != '/') and (p1[0] == '/')
        return linux_path(p0+p1)

def mk_key(xy):
    x,_ = xy
    begin, end  =  '_Ktrials_', '.tec'
    assert x.count(begin) == x.count(end) == 1
    base, trials = x.split(begin)
    trials       = lmap(int,trials.rstrip(end).split('_'))
    ind          = int(basename(x).split('_ind_')[1].split('_')[0])
    return ind, base, len(trials), trials, xy

def add_to_dict(D, val, keys):
    for key in keys[:-1]:
        if not key in D:
            D[key] = {}
        D = D[key]
    D[keys[-1]] = val

def regroup_arch_pics(A, isvalid = None):
    result = {}
    for D in A:
        if isvalid(D):
            add_to_dict(result,  D, [D['info_tec']['tag_run']  ,
                                     D['info_tec']['qmode']    ,
                                     D['info_tec']['tec_fname'],
                                     D['key_scalar']           ])
    return result

if __name__ == '__main__':
    
    for     tag_run  in ['_rnd12_']:
        for qmode    in ['Nu_classic']:
            archive_pics_full  =  regroup_arch_pics(eval(' '.join(reader(pjoin(post_output,f'pics_archive_tag_run{tag_run}qmode_{qmode}.dat')))),
                                                    isvalid = lambda D: len(D['info_tec']['K_trials']) == 2                                                    )

            archive_pics       =  sorted(archive_pics_full[tag_run][qmode].items(),key = mk_key)
            fname_pptx         =  pjoin(post_output,'presentations',f'{tstamp()}_{qmode}_tag_run_{tag_run}.pptx')
            D_qtypes           =  {'Nu_classic': {'qmode':'Nu_classic'}}[qmode]
            assert all(map(lambda x: qmode in x[0], archive_pics))

            all_fname_pics  =  []

            for k,v in archive_pics:
                all_fname_pics.append(dict(v.items()))
                all_fname_pics[-1]['slide_name'] = basename(k).replace('.tec','')

            Mk_Presentation(all_fname_pics, fname_pptx, D_qtypes)

